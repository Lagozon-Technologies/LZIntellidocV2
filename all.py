import os
import streamlit as st
from langchain.text_splitter import CharacterTextSplitter
from langchain_openai import OpenAIEmbeddings
from langchain_community.vectorstores import FAISS
from PIL import Image
import pandas as pd
from PyPDF2 import PdfReader 
from docx import Document
from pptx import Presentation
import pdfplumber  #for identifying page number -- sujal 12/07/24
import easyocr    #for scan document  -- sujal 12/07/24
import numpy as np
from openai import OpenAI
import hashlib   #for tracking  -- sujal 12/07/24

# Uncomment the following line if you face the OpenMP error mentioned in the comments
os.environ["KMP_DUPLICATE_LIB_OK"] = "TRUE"

api_key = st.secrets["API_KEY"]

client = OpenAI(api_key=api_key)

img = Image.open(r"images.png")
st.set_page_config(page_title="DocGenius: Document Generation AI", page_icon=img)


st.sidebar.title("IntelliDoc")
# Department dropdown in the sidebar
department = st.sidebar.selectbox("Select Department", ["HR", "Legal", "L&D", "Operations", "Insurance"])


# Select file type in the sidebar
file_type = st.sidebar.selectbox("Select file type", ["PDF", "Scanned PDF", "DOCX", "PPTX", "EXCEL", "CSV"])
accepted_types = {"PDF": "pdf", "Scanned PDF": "pdf", "DOCX": "docx", "PPTX": "pptx", "EXCEL": "xlsx", "CSV": "csv"}
file_extension = accepted_types[file_type]



col1, col2 = st.columns([1, 4])

with col1:
    st.image("lagozon_technologies_private_limited_logo.jpeg", width=90)

with col2:
    st.title("Ask Your Documents 📄")


if "previous_file_type" not in st.session_state:
    st.session_state.previous_file_type = file_type

if st.session_state.previous_file_type != file_type:
    st.session_state.previous_file_type = file_type
    if "conversation" in st.session_state:
        del st.session_state.conversation

uploaded_files = st.sidebar.file_uploader(f"Upload your {file_type} files", type=file_extension, accept_multiple_files=True, on_change=lambda: reset_session())

# Clear session button in the sidebar
if st.sidebar.button("Clear"):
    if "conversation" in st.session_state:
        del st.session_state.conversation
    if "previous_files_hash" in st.session_state:
        del st.session_state.previous_files_hash
    st.session_state.query = ""
    st.experimental_rerun()
    


def reset_session():
    if "conversation" in st.session_state:
        del st.session_state.conversation
    if "previous_files_hash" in st.session_state:
        del st.session_state.previous_files_hash
    st.session_state.query = ""

def get_file_hash(file):
    file.seek(0)
    file_bytes = file.read()
    file.seek(0)
    return hashlib.md5(file_bytes).hexdigest()

def extract_text_from_scanned_pdf(file):
    reader = easyocr.Reader(['en'])
    text = ""

    with pdfplumber.open(file) as pdf:
        for page in pdf.pages:
            page_image = page.to_image()
            image_np = np.array(page_image.original)
            result = reader.readtext(image_np, detail=0, paragraph=True)
            text += " ".join(result)

    return text

if uploaded_files:
    current_files_hash = [get_file_hash(file) for file in uploaded_files]
    if "previous_files_hash" not in st.session_state:
        st.session_state.previous_files_hash = current_files_hash

    if st.session_state.previous_files_hash != current_files_hash:
        st.session_state.previous_files_hash = current_files_hash
        if "conversation" in st.session_state:
            del st.session_state.conversation
        st.session_state.query = ""

    all_texts = []
    for uploaded_file in uploaded_files:
        if file_type == "PDF":
            pdf_reader = PdfReader(uploaded_file)
            text = "".join([page.extract_text() or "" for page in pdf_reader.pages])
            if text:
                all_texts.append(text)
        elif file_type == "Scanned PDF":
            text = extract_text_from_scanned_pdf(uploaded_file)
            all_texts.append(text)
        elif file_type == "DOCX":
            doc = Document(uploaded_file)
            doc_text = '\n'.join([paragraph.text for paragraph in doc.paragraphs])
            all_texts.append(doc_text)
        elif file_type == "PPTX":
            ppt = Presentation(uploaded_file)
            ppt_text = '\n'.join([shape.text for slide in ppt.slides for shape in slide.shapes if hasattr(shape, "text")])
            all_texts.append(ppt_text)
        elif file_type == "EXCEL":
            df = pd.read_excel(uploaded_file)
            excel_text = df.to_string(index=False)
            all_texts.append(excel_text)
        elif file_type == "CSV":
            df = pd.read_csv(uploaded_file)
            csv_text = df.to_string(index=False)
            all_texts.append(csv_text)

    if all_texts:
        text_splitter = CharacterTextSplitter(separator="\n", chunk_size=1000, chunk_overlap=200, length_function=len)
        chunks = []
        for text in all_texts:
            chunks.extend(text_splitter.split_text(text)) 

        embeddings = OpenAIEmbeddings(api_key=api_key)
        knowledge_base = FAISS.from_texts(chunks, embeddings)

        llm_model = st.selectbox("Select LLM Model", ["gpt-3.5-turbo", "gpt-4", "gpt-4-turbo"])

        query = st.text_input("Ask your Question about your documents", value=st.session_state.get("query", ""), key="query")
        if "conversation" not in st.session_state:
            st.session_state.conversation = []

        if query:
            docs = knowledge_base.similarity_search(query)
            doc_texts = "\n".join([doc.page_content for doc in docs])
            prompt = f"Answer the question based on the following documents:\n\n{doc_texts}\n\nQuestion: {query}"

            response = client.chat.completions.create(
                messages=[
                    {"role": "system", "content": "You are a helpful assistant."},
                    {"role": "user", "content": prompt}
                ],
                model=llm_model
            )
            answer = response.choices[0].message.content
            st.session_state.conversation.append((query, answer))

        for question, answer in st.session_state.conversation:
            st.write(f"*Question:* {question}")
            st.write(f"*Answer:* {answer}")